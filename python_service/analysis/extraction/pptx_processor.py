import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_shape(shape):
    text_runs = []
    
    # 1. 普通文本框
    if hasattr(shape, "text"):
        text_runs.append(shape.text)
    
    # 2. 组合形状：递归提取
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text_runs.extend(extract_text_from_shape(sub_shape))
    
    # 3. 表格：逐行逐列提取
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            text_runs.append("\t".join(row_text))
    
    return text_runs
def batch_extract_pptx(folder_path, output_folder="../data/processed"):
    os.makedirs(output_folder, exist_ok=True)
    
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx") and not filename.startswith("~$"):  # 跳过临时文件
            pptx_path = os.path.join(folder_path, filename)
            print(f"正在处理: {filename}")
            
            try:
                text = extract_text_from_pptx_advanced(pptx_path)
                output_path = os.path.join(output_folder, f"{os.path.splitext(filename)[0]}.txt")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
            except Exception as e:
                print(f"处理失败 {filename}: {e}")
def extract_text_from_pptx_advanced(pptx_path):
    prs = Presentation(pptx_path)
    full_text = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        full_text.append(f"\n--- 第 {slide_idx} 页 ---\n")
        for shape in slide.shapes:
            full_text.extend(extract_text_from_shape(shape))
    
    return "\n".join(full_text)
# 使用示例：把你的 PPTX 文件夹路径填进去
batch_extract_pptx("../data/raw")